import tempfile
import os
import fitz  # PyMuPDF
import camelot
import numpy as np
from PIL import Image
import io
from PyPDF2 import PdfReader
import pdfplumber
import contextlib
import pandas as pd
import time
import docx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

class PDFExtractor:
    def __init__(self, file_path: str, page_range: tuple = None):
        self.file_path = file_path
        self.page_range = page_range
        self._pdf_document = None
        
    @property
    def pdf_document(self):
        """Lazy loading of PDF document"""
        if self._pdf_document is None:
            self._pdf_document = fitz.open(self.file_path)
        return self._pdf_document
        
    def __enter__(self):
        return self
        
    def __exit__(self, exc_type, exc_val, exc_tb):
        self.close()
        
    def close(self):
        """Explicitly close the PDF document"""
        if hasattr(self, '_pdf_document') and self._pdf_document:
            try:
                self._pdf_document.close()
                self._pdf_document = None
            except Exception:
                pass
                
    def __del__(self):
        """Backup method to ensure resources are freed"""
        self.close()
        
    def _get_page_range(self, total_pages):
        """Get normalized page range with 0-based indexing"""
        if self.page_range and len(self.page_range) >= 2:
            start_page = max(0, self.page_range[0] - 1)  # Convert to 0-based index
            end_page = min(total_pages, self.page_range[1])
            return start_page, end_page
        else:
            return 0, total_pages
    
    def _get_page_range_1_based(self, total_pages):
        """Get normalized page range with 1-based indexing (for camelot)"""
        if self.page_range and len(self.page_range) >= 2:
            start_page = max(1, self.page_range[0])  # Keep as 1-based index
            end_page = min(total_pages, self.page_range[1])
            return start_page, end_page
        else:
            return 1, total_pages
        
    def extract_text(self):
        """Extract text content using PyPDF2"""
        text_content = []
        
        try:
            reader = PdfReader(self.file_path)
            text = ""
            total_pages = len(reader.pages)
            start_page, end_page = self._get_page_range(total_pages)
            
            for i, page in enumerate(reader.pages):
                if start_page <= i < end_page:
                    extracted_text = page.extract_text()
                    if extracted_text:  # Only add non-empty text
                        text += extracted_text + "\n"
            
            # Only add if we have actual content
            if text.strip():
                text_content.append({
                    'content': text,
                    'metadata': {
                        'source': self.file_path,
                        'page_range': self.page_range
                    }
                })
        except Exception as e:
            raise Exception(f"Error extracting text content: {str(e)}")
        
        return text_content
        
    def extract_images(self):
        """Extract images using PyMuPDF"""
        image_content = []
        
        try:
            total_pages = len(self.pdf_document)
            start_page, end_page = self._get_page_range(total_pages)
            
            for page_num in range(start_page, end_page):
                page = self.pdf_document[page_num]
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = self.pdf_document.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # Convert to PIL Image for validation and processing
                        with Image.open(io.BytesIO(image_bytes)) as image:
                            # Skip tiny images (often artifacts)
                            if image.width < 50 or image.height < 50:
                                continue
                                
                            # Convert PIL Image to bytes
                            img_byte_arr = io.BytesIO()
                            image.save(img_byte_arr, format=image.format if image.format else 'PNG')
                            img_format = image.format if image.format else 'PNG'
                            img_size = image.size
                            
                        image_content.append({
                            'content': img_byte_arr.getvalue(),
                            'metadata': {
                                'page_number': page_num + 1,
                                'image_index': img_index,
                                'format': img_format,
                                'size': img_size
                            }
                        })
                    except Exception as e:
                        print(f"Warning: Error extracting image {img_index} on page {page_num + 1}: {str(e)}")
        except Exception as e:
            raise Exception(f"Error extracting images: {str(e)}")
            
        return image_content
        
    def extract_tables(self, table_flavor="auto", debug=True):
        """Extract tables using the specified method with enhanced debugging"""
        table_content = []
        debug_stats = {
            "lattice": {"found": 0, "filtered": 0, "reasons": []},
            "stream": {"found": 0, "filtered": 0, "reasons": []}, 
            "pdfplumber": {"found": 0, "filtered": 0, "reasons": []}
        }

        try:
            # Get page range with 1-based indexing for camelot
            with fitz.open(self.file_path) as pdf:
                total_pages = pdf.page_count
                if debug:
                    print(f"Total pages: {total_pages}")
            
            start_page, end_page = self._get_page_range_1_based(total_pages)
            page_str = ','.join(str(i) for i in range(start_page, end_page + 1))
            
            # Only use camelot-lattice if selected or auto
            if table_flavor in ["auto", "lattice"]:
                try:
                    tables_lattice = camelot.read_pdf(
                        self.file_path, 
                        pages=page_str, 
                        flavor='lattice',
                        suppress_stdout=not debug
                    )
                    
                    if debug:
                        print(f"Camelot lattice found {len(tables_lattice)} potential tables")
                    
                    debug_stats["lattice"]["found"] = len(tables_lattice)
                    
                    for i, table in enumerate(tables_lattice):
                        # More relaxed filtering
                        if len(table.df) > 0:  # Just check if there's any data
                            df = table.df.replace('', np.nan).dropna(how='all').fillna('')
                            
                            # Simplified filtering - only filter empty tables
                            if not df.empty and df.shape[0] > 0 and df.shape[1] > 0:
                                table_content.append({
                                    'content': df.to_dict(orient='records'),
                                    'metadata': {
                                        'page_number': table.page,
                                        'table_index': i,
                                        'accuracy': table.accuracy,
                                        'method': 'camelot-lattice',
                                        'shape': df.shape
                                    }
                                })
                            else:
                                debug_stats["lattice"]["filtered"] += 1
                                debug_stats["lattice"]["reasons"].append(f"Empty table after cleaning: {df.shape}")
                        else:
                            debug_stats["lattice"]["filtered"] += 1
                            debug_stats["lattice"]["reasons"].append("Empty table")
                except Exception as e:
                    if debug:
                        print(f"Warning: Camelot lattice extraction error: {str(e)}")
            
            # Only use camelot-stream if selected or auto
            if table_flavor in ["auto", "stream"]:
                try:
                    tables_stream = camelot.read_pdf(
                        self.file_path, 
                        pages=page_str, 
                        flavor='stream',
                        suppress_stdout=not debug
                    )
                    
                    if debug:
                        print(f"Camelot stream found {len(tables_stream)} potential tables")
                    
                    debug_stats["stream"]["found"] = len(tables_stream)
                    
                    for i, table in enumerate(tables_stream):
                        # Relaxed filtering - lower accuracy threshold
                        if table.accuracy > 50 and len(table.df) > 0:  # Lower threshold from 90 to 50
                            df = table.df.replace('', np.nan).dropna(how='all').fillna('')
                            
                            # Simplified checks
                            if not df.empty and df.shape[0] >= 2 and df.shape[1] >= 2:
                                # Simpler duplicate check
                                is_duplicate = False
                                for existing in table_content:
                                    if existing['metadata']['page_number'] == table.page:
                                        if existing['metadata']['method'] == 'camelot-lattice':
                                            # Compare shapes to detect potential duplicates
                                            existing_shape = existing['metadata'].get('shape', (0, 0))
                                            if existing_shape == df.shape:
                                                is_duplicate = True
                                                break
                                
                                if not is_duplicate:
                                    # Removed the avg_word_count filtering
                                    table_content.append({
                                        'content': df.to_dict(orient='records'),
                                        'metadata': {
                                            'page_number': table.page,
                                            'table_index': i,
                                            'accuracy': table.accuracy,
                                            'method': 'camelot-stream',
                                            'shape': df.shape
                                        }
                                    })
                                else:
                                    debug_stats["stream"]["filtered"] += 1
                                    debug_stats["stream"]["reasons"].append("Duplicate of lattice table")
                            else:
                                debug_stats["stream"]["filtered"] += 1
                                debug_stats["stream"]["reasons"].append(f"Table too small after cleaning: {df.shape}")
                        else:
                            debug_stats["stream"]["filtered"] += 1
                            if table.accuracy <= 50:
                                debug_stats["stream"]["reasons"].append(f"Low accuracy: {table.accuracy}")
                            else:
                                debug_stats["stream"]["reasons"].append("Empty table")
                except Exception as e:
                    if debug:
                        print(f"Warning: Camelot stream extraction error: {str(e)}")
            
            # Only use pdfplumber if selected or auto
            if table_flavor in ["auto", "pdfplumber"]:
                try:
                    with pdfplumber.open(self.file_path) as pdf:
                        start_page, end_page = self._get_page_range(len(pdf.pages))  # 0-based for pdfplumber
                        
                        total_found = 0
                        for page_num in range(start_page, end_page):
                            page = pdf.pages[page_num]
                            
                            # Try multiple table extraction methods
                            table_settings = [
                                # Regular lines strategy
                                {"vertical_strategy": "lines", "horizontal_strategy": "lines"},
                                # Text-based strategy for tables without explicit lines
                                {"vertical_strategy": "text", "horizontal_strategy": "text"}
                            ]
                            
                            for settings in table_settings:
                                try:
                                    tables = page.extract_tables(table_settings=settings)
                                    total_found += len(tables)
                                    
                                    for idx, table in enumerate(tables):
                                        if table and len(table) >= 2:  # Reduced from 3 to 2
                                            rows = []
                                            for row in table:
                                                # Include rows with at least one non-empty cell
                                                if row and any(cell and str(cell).strip() for cell in row):
                                                    rows.append(row)
                                            
                                            if rows and len(rows) >= 2:  # Reduced from 3 to 2
                                                df = pd.DataFrame(rows)
                                                
                                                # Use first row as header if it contains text
                                                if all(isinstance(x, str) and x.strip() for x in df.iloc[0] if x is not None):
                                                    df.columns = df.iloc[0]
                                                    df = df.iloc[1:].reset_index(drop=True)
                                                
                                                df = df.replace('', np.nan).dropna(how='all', axis=1).dropna(how='all').fillna('')
                                                
                                                # Relaxed filtering - only basic size check
                                                if not df.empty and df.shape[0] >= 1 and df.shape[1] >= 2:
                                                    table_content.append({
                                                        'content': df.to_dict(orient='records'),
                                                        'metadata': {
                                                            'page_number': page_num + 1,
                                                            'table_index': idx,
                                                            'method': f'pdfplumber-{settings["vertical_strategy"]}',
                                                            'shape': df.shape
                                                        }
                                                    })
                                                else:
                                                    debug_stats["pdfplumber"]["filtered"] += 1
                                                    debug_stats["pdfplumber"]["reasons"].append(f"Table too small after cleaning: {df.shape}")
                                            else:
                                                debug_stats["pdfplumber"]["filtered"] += 1
                                                debug_stats["pdfplumber"]["reasons"].append(f"Not enough rows: {len(rows)}")
                                        else:
                                            debug_stats["pdfplumber"]["filtered"] += 1
                                            debug_stats["pdfplumber"]["reasons"].append("Not enough rows in raw table")
                                except Exception as e:
                                    if debug:
                                        print(f"Warning: pdfplumber error with settings {settings}: {str(e)}")
                    
                        debug_stats["pdfplumber"]["found"] = total_found
                        if debug:
                            print(f"PDFPlumber found {total_found} potential tables")
                except Exception as e:
                    if debug:
                        print(f"Warning: pdfplumber error: {str(e)}")
        
            # Print debug summary
            if debug:
                print("\n=== Table Extraction Summary ===")
                print(f"Total tables extracted: {len(table_content)}")
                for method, stats in debug_stats.items():
                    print(f"\n{method.upper()}")
                    print(f"  Found: {stats['found']}")
                    print(f"  Filtered out: {stats['filtered']}")
                    if stats['reasons']:
                        print(f"  Reasons (sample): {stats['reasons'][:5]}")
                print("==============================\n")
                
        except Exception as e:
            raise Exception(f"Error extracting tables: {str(e)}")

        return table_content

@contextlib.contextmanager
def temporary_file(content, suffix=None):
    """Simple context manager for temporary files"""
    temp_file = None
    try:
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        temp_file.write(content)
        temp_file.close()
        yield temp_file.name
    finally:
        if temp_file and os.path.exists(temp_file.name):
            try:
                os.remove(temp_file.name)
            except Exception:
                pass

def process_pdf(file_path, page_range=None, extract_text=True, extract_images=True, 
                extract_tables=True, table_flavor="auto", debug=False):
    """Process PDF file and return extracted content with debugging capability"""
    # Validate PDF file before processing
    try:
        with open(file_path, 'rb') as f:
            if not f.read(4).startswith(b'%PDF'):
                raise ValueError("The file does not appear to be a valid PDF.")
    except Exception as e:
        raise Exception(f"Error validating PDF file: {str(e)}")
    
    content = {'text': [], 'images': [], 'tables': [], 'debug_info': {}}
    
    if debug:
        print(f"Processing PDF: {file_path}")
        print(f"Page range: {page_range if page_range else 'all pages'}")
        print(f"Extraction flags: text={extract_text}, images={extract_images}, tables={extract_tables}")
        print(f"Table flavor: {table_flavor}")
    
    with PDFExtractor(file_path, page_range) as extractor:
        if extract_text:
            start_time = time.time()
            content['text'] = extractor.extract_text()
            if debug:
                content['debug_info']['text_extraction_time'] = time.time() - start_time
                print(f"Text extraction complete. Found {len(content['text'])} text blocks.")
        
        if extract_images:
            start_time = time.time()
            content['images'] = extractor.extract_images()
            if debug:
                content['debug_info']['image_extraction_time'] = time.time() - start_time
                print(f"Image extraction complete. Found {len(content['images'])} images.")
        
        if extract_tables:
            start_time = time.time()
            content['tables'] = extractor.extract_tables(table_flavor=table_flavor, debug=debug)
            if debug:
                content['debug_info']['table_extraction_time'] = time.time() - start_time
                print(f"Table extraction complete. Found {len(content['tables'])} tables.")
    
    return content

def get_pdf_metadata(file_path):
    """Extract metadata from PDF file"""
    metadata = {}
    
    try:
        with fitz.open(file_path) as pdf:
            metadata = {
                'title': pdf.metadata.get('title', ''),
                'author': pdf.metadata.get('author', ''),
                'subject': pdf.metadata.get('subject', ''),
                'keywords': pdf.metadata.get('keywords', ''),
                'creator': pdf.metadata.get('creator', ''),
                'producer': pdf.metadata.get('producer', ''),
                'creationDate': pdf.metadata.get('creationDate', ''),
                'modDate': pdf.metadata.get('modDate', ''),
                'total_pages': pdf.page_count,
                'file_size': os.path.getsize(file_path),
                'has_toc': bool(pdf.get_toc())
            }
    except Exception as e:
        print(f"Warning: Error extracting PDF metadata: {str(e)}")
    
    return metadata

def get_table_summary(tables):
    """Provide a summary of extracted tables"""
    summary = {
        'total_tables': len(tables),
        'tables_by_method': {},
        'tables_by_page': {},
        'average_rows': 0,
        'average_cols': 0
    }
    
    if not tables:
        return summary
    
    total_rows = 0
    total_cols = 0
    
    for table in tables:
        method = table['metadata'].get('method', 'unknown')
        page = table['metadata'].get('page_number', 0)
        
        # Count by method
        if method not in summary['tables_by_method']:
            summary['tables_by_method'][method] = 0
        summary['tables_by_method'][method] += 1
        
        # Count by page
        if page not in summary['tables_by_page']:
            summary['tables_by_page'][page] = 0
        summary['tables_by_page'][page] += 1
        
        # Calculate table dimensions
        if table['content']:
            rows = len(table['content'])
            cols = len(table['content'][0].keys()) if rows > 0 else 0
            total_rows += rows
            total_cols += cols
    
    # Calculate averages
    if summary['total_tables'] > 0:
        summary['average_rows'] = total_rows / summary['total_tables']
        summary['average_cols'] = total_cols / summary['total_tables']
    
    return summary

def process_pdfs_batch(file_paths, **kwargs):
    """Process multiple PDF files with the same extraction parameters"""
    results = {}
    errors = {}
    
    for file_path in file_paths:
        try:
            results[file_path] = process_pdf(file_path, **kwargs)
        except Exception as e:
            errors[file_path] = str(e)
            print(f"Error processing {file_path}: {str(e)}")
    
    return {
        'results': results,
        'errors': errors,
        'success_rate': len(results) / len(file_paths) if file_paths else 0
    }

def process_doc(file_path):
    """Process DOC file and return extracted content"""
    content = {'text': [], 'tables': [], 'images': []}
    
    try:
        doc = docx.Document(file_path)
        
        # Extract text
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text.strip():
            content['text'].append({
                'content': text,
                'metadata': {
                    'source': file_path
                }
            })
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            content['tables'].append({
                'content': table_data,
                'metadata': {
                    'source': file_path
                }
            })
        
        # Extract images
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                image = rel.target_part.blob
                content['images'].append({
                    'content': image,
                    'metadata': {
                        'source': file_path
                    }
                })
    except Exception as e:
        raise Exception(f"Error processing DOC file: {str(e)}")
    
    return content

def process_ppt(file_path):
    """Process PPT file and return extracted content"""
    content = {'text': [], 'tables': [], 'images': []}
    
    try:
        ppt = Presentation(file_path)
        
        # Extract text
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        if text.strip():
            content['text'].append({
                'content': text,
                'metadata': {
                    'source': file_path
                }
            })
        
        # Extract tables
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text_frame.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    content['tables'].append({
                        'content': table_data,
                        'metadata': {
                            'source': file_path
                        }
                    })
        
        # Extract images
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image.blob
                    content['images'].append({
                        'content': image,
                        'metadata': {
                            'source': file_path
                        }
                    })
    except Exception as e:
        raise Exception(f"Error processing PPT file: {str(e)}")
    
    return content

def process_file(file_path, **kwargs):
    """Process file based on its type and return extracted content"""
    if file_path.lower().endswith('.pdf'):
        return process_pdf(file_path, **kwargs)
    elif file_path.lower().endswith('.doc') or file_path.lower().endswith('.docx'):
        return process_doc(file_path)
    elif file_path.lower().endswith('.ppt') or file_path.lower().endswith('.pptx'):
        return process_ppt(file_path)
    else:
        raise ValueError("Unsupported file type")

def process_files_batch(file_paths, **kwargs):
    """Process multiple files with the same extraction parameters"""
    results = {}
    errors = {}
    
    for file_path in file_paths:
        try:
            results[file_path] = process_file(file_path, **kwargs)
        except Exception as e:
            errors[file_path] = str(e)
            print(f"Error processing {file_path}: {str(e)}")
    
    return {
        'results': results,
        'errors': errors,
        'success_rate': len(results) / len(file_paths) if file_paths else 0
    }