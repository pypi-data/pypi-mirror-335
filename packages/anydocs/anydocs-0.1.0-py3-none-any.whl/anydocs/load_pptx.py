import typing as tp
from dataclasses import dataclass

import base64c as base64  # type: ignore
from pptx import Presentation

from ._base import Artifact


@dataclass
class PptxLoader(Artifact):
    def extract_text(self) -> tp.Generator[str, None, None]:
        prs = Presentation(self.file_path)
        for slide in prs.slides:  # type: ignore
            for shape in slide.shapes:  # type: ignore
                if shape.has_text_frame:  # type: ignore
                    text_frame = shape.text_frame  # type: ignore
                    for paragraph in text_frame.paragraphs:  # type: ignore
                        if paragraph.text:  # type: ignore
                            yield paragraph.text  # type: ignore
                        else:
                            continue

    def extract_image(self):
        prs = Presentation(self.file_path)
        for slide in prs.slides:  # type: ignore
            for shape in slide.shapes:  # type: ignore
                if shape.shape_type == 13:  # type: ignore
                    image = shape.image  # type: ignore
                    yield base64.b64encode(image.blob).decode()
                else:
                    continue
