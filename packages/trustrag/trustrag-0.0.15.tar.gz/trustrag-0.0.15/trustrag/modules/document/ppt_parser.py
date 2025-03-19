from io import BytesIO
from pptx import Presentation


class PptParser(object):
    def __init__(self):
        super().__init__()

    def __extract(self, shape):
        if shape.shape_type == 19:
            tb = shape.table
            rows = []
            for i in range(1, len(tb.rows)):
                rows.append("; ".join([tb.cell(
                    0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
            return "\n".join(rows)

        if shape.has_text_frame:
            return shape.text_frame.text

        if shape.shape_type == 6:
            texts = []
            for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                t = self.__extract(p)
                if t:
                    texts.append(t)
            return "\n".join(texts)

    def parse(self, fnm, from_page=0, to_page=100000, callback=None):
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in sorted(
                    slide.shapes, key=lambda x: (x.top // 10, x.left)):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))

        return txts

if __name__ == '__main__':
    pp=PptParser()
    txts=pp.parse('/data/users/searchgpt/yq/GoMate/data/docs/集成电路封装材料-光敏材料.pptx')
    print(txts)
